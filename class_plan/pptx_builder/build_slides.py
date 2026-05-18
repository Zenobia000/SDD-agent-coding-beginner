"""Vibe Coding 一日工作坊 — 初版 PPT 產生器（Figma 風格：白底、黑字、Typography 主導）.

依據 class_plan/slides_outline.md 的 59 張投影片 + 附錄三個備選 Demo 專案（共 64 張）。
產物：class_plan/vibe_coding_workshop.pptx

設計哲學（對齊 .claude/ui/figma/DESIGN.md）：
- 介面層純黑白：所有文字/框線/分隔都是 #111 或灰階
- 色彩只出現在 hero 段落（封面、模組分隔、金句）的薄漸層條
- 大量留白（Gallery pacing）
- Typography 做層級：用字級差拉層次，不靠顏色
- Mono 小寫做段落標籤（M0 · OPENING）
- Pill 幾何（圓角膠囊）代替矩形
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# --- 設計 Token（Figma-inspired, 極簡黑白）--------------------------------
BG = RGBColor(0xFF, 0xFF, 0xFF)  # 純白
BG_SOFT = RGBColor(0xFA, 0xFA, 0xF7)  # 實作頁：極淡暖白
INK = RGBColor(0x11, 0x11, 0x11)  # 近黑（比純黑柔和）
INK_SOFT = RGBColor(0x3A, 0x3A, 0x3A)  # 次級文字
MUTED = RGBColor(0x8A, 0x8A, 0x8A)  # 淡灰（說明、頁碼）
HAIRLINE = RGBColor(0xE5, 0xE5, 0xE5)  # 極淡分隔線
CARD_BORDER = RGBColor(0xE8, 0xE8, 0xE5)  # 卡片邊框

# 單一強調色（用得極省，只給「重要動詞」用）
ACCENT = RGBColor(0x11, 0x11, 0x11)  # 強調用純黑
ACCENT_WARN = RGBColor(0xB4, 0x53, 0x09)  # 暖褐（❌ 壞範例）
ACCENT_GOOD = RGBColor(0x15, 0x6B, 0x4A)  # 深綠（✅ 好範例）
HIGHLIGHT_BG = RGBColor(0xFE, 0xF7, 0xE6)  # 淡黃底（Highlight）

# Hero 漸層用色（只給封面用薄色塊）
HERO_PINK = RGBColor(0xF9, 0xD5, 0xE5)
HERO_PURPLE = RGBColor(0xDD, 0xCC, 0xF0)
HERO_YELLOW = RGBColor(0xFD, 0xF0, 0xCE)
HERO_GREEN = RGBColor(0xCC, 0xE6, 0xD3)

# 字體：Windows 原生中文 + Consolas
FONT_TITLE = "Microsoft JhengHei Light"
FONT_BODY = "Microsoft JhengHei"
FONT_BOLD = "Microsoft JhengHei"
FONT_MONO = "Consolas"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_PATH = Path(__file__).resolve().parent.parent / "vibe_coding_workshop.pptx"


# --- Helper ---------------------------------------------------------------


def _new_slide(prs: Presentation, bg: RGBColor = BG):
    """新增一張空白投影片並填滿底色."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.shadow.inherit = False
    spTree = bg_shape._element.getparent()
    spTree.remove(bg_shape._element)
    spTree.insert(2, bg_shape._element)
    return slide


def _add_text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    size: int = 20,
    bold: bool = False,
    color: RGBColor = INK,
    font: str = FONT_BODY,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    line_spacing: float = 1.35,
    letter_spacing: float = 0.0,
):
    """在指定框加入文字（單段）。letter_spacing 單位 pt（負值收緊）."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        # 字元間距（負值收緊，對應 Figma 的 negative tracking）
        if letter_spacing != 0:
            from pptx.oxml.ns import qn
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(letter_spacing * 100)))
    return tx


def _hairline(slide, left, top, width, *, color: RGBColor = HAIRLINE, weight: float = 0.5):
    """水平細線."""
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def _add_footer(slide, module_label: str, slide_no: int, total: int = 64):
    """頁尾：左下模組標籤（Mono Caps）、右下頁碼."""
    _hairline(slide, Inches(0.6), Inches(7.05), Inches(12.13))
    _add_text(
        slide,
        module_label.upper(),
        Inches(0.6),
        Inches(7.15),
        Inches(8),
        Inches(0.3),
        size=10,
        color=MUTED,
        font=FONT_MONO,
        letter_spacing=2.0,
    )
    _add_text(
        slide,
        f"{slide_no:02d} / {total:02d}",
        Inches(11.73),
        Inches(7.15),
        Inches(1),
        Inches(0.3),
        size=10,
        color=MUTED,
        font=FONT_MONO,
        align=PP_ALIGN.RIGHT,
        letter_spacing=2.0,
    )


def _mono_label(slide, text: str, left, top, width=Inches(8), *, color: RGBColor = MUTED, size: int = 11):
    """Mono caps 段落標籤."""
    _add_text(
        slide,
        text.upper(),
        left,
        top,
        width,
        Inches(0.3),
        size=size,
        color=color,
        font=FONT_MONO,
        letter_spacing=2.5,
    )


def _pill(slide, text: str, left, top, width=Inches(1.6), height=Inches(0.38), *,
          fill: RGBColor = INK, text_color: RGBColor = BG, font_size: int = 11, mono: bool = True):
    """Pill-shaped chip（Figma 招牌）."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # 圓角調整到極大 → pill 效果
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper() if mono else text
    r.font.name = FONT_MONO if mono else FONT_BOLD
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = text_color
    if mono:
        from pptx.oxml.ns import qn
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", "200")
    return shape


def _card(slide, left, top, width, height, *, bg: RGBColor = BG, border: RGBColor = CARD_BORDER):
    """白卡 + 細邊框，無陰影."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.03  # 8px 左右圓角
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    return shape


def _add_code_block(slide, code: str, left, top, width, height, *, size: int = 13):
    """碼塊：極淡底、細邊框、mono 字."""
    shape = _card(slide, left, top, width, height, bg=RGBColor(0xFA, 0xFA, 0xFA), border=HAIRLINE)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    for idx, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.3
        r = p.add_run()
        r.text = line
        r.font.name = FONT_MONO
        r.font.size = Pt(size)
        r.font.color.rgb = INK
    return shape


def _hero_gradient_strip(slide, top=Inches(0), height=Inches(0.35)):
    """封面用的薄彩條（Figma 招牌漸層精神，但壓成細條避免壓迫）."""
    colors = [HERO_PINK, HERO_PURPLE, HERO_YELLOW, HERO_GREEN]
    w = SLIDE_W / len(colors)
    for i, c in enumerate(colors):
        block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(w * i)), top, Emu(int(w)), height)
        block.fill.solid()
        block.fill.fore_color.rgb = c
        block.line.fill.background()


# --- Slide Builders --------------------------------------------------------


def s_cover(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.3))
    _hero_gradient_strip(s, top=Inches(7.2), height=Inches(0.3))

    _mono_label(s, "Vibe Coding · One-day Workshop", Inches(0.8), Inches(1.2))
    _add_text(
        s, "用三個免費工具",
        Inches(0.8), Inches(1.85),
        Inches(11.7), Inches(1.2),
        size=68, bold=True, color=INK, font=FONT_BOLD, letter_spacing=-3.0,
    )
    _add_text(
        s, "把你腦中的想法做成可用的小專案",
        Inches(0.8), Inches(2.9),
        Inches(11.7), Inches(1.2),
        size=68, bold=True, color=INK, font=FONT_BOLD, letter_spacing=-3.0,
    )

    _hairline(s, Inches(0.8), Inches(4.7), Inches(11.7))

    # 三工具
    tools = ["Google AI Studio", "Antigravity", "Gemini CLI"]
    for i, t in enumerate(tools):
        left = Inches(0.8 + i * 4.1)
        _mono_label(s, f"Stage {i + 1:02d}", left, Inches(5.0), width=Inches(4))
        _add_text(s, t, left, Inches(5.35), Inches(4), Inches(0.6), size=22, bold=True, letter_spacing=-0.5)

    _add_text(s, "[ 講師 ] · [ 日期 ]", Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4), size=12, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
    return s


def s_you_too(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M0 · The honest question", Inches(0.8), Inches(0.7))
    _add_text(s, "你是不是也這樣？", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=44, bold=True, letter_spacing=-1.5)

    quotes = [
        '「ChatGPT 問過、Gemini 問過，但我還是不會做出任何東西。」',
        '「我有很多想法，但不會寫程式，也不想學。」',
        '「聽說 Vibe Coding 很紅，但那是工程師的事吧？」',
    ]
    for i, q in enumerate(quotes):
        top = Inches(2.8 + i * 1.15)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, q, Inches(0.8), top + Inches(0.25), Inches(11.7), Inches(0.9),
                  size=22, color=INK_SOFT, letter_spacing=-0.3)
    _hairline(s, Inches(0.8), Inches(6.25), Inches(11.7))

    _add_footer(s, "M0 · Opening", no)
    return s


def s_promise(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.3))
    _mono_label(s, "Today's promise", Inches(0.8), Inches(1.2))
    _add_text(s, "今天結束前，你會做出一個", Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8),
              size=26, color=MUTED, letter_spacing=-0.3)

    _add_text(s, "可以打開來用", Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2),
              size=90, bold=True, letter_spacing=-3.5)
    _add_text(s, "可以秀給朋友看的", Inches(0.8), Inches(3.75), Inches(11.7), Inches(1.2),
              size=90, bold=True, letter_spacing=-3.5)
    _add_text(s, "小  專  案。", Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.2),
              size=90, bold=True, letter_spacing=-3.5, color=INK)

    _hairline(s, Inches(0.8), Inches(6.3), Inches(11.7))
    _add_footer(s, "M0 · Opening", no)
    return s


def s_three_tools(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M0 · The journey", Inches(0.8), Inches(0.7))
    _add_text(s, "三個工具  ·  一條線", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=44, bold=True, letter_spacing=-1.5)
    _add_text(s, "同一個專案，三個工具輪流接棒。", Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2)

    titles = ["早上", "下午", "傍晚"]
    names = ["AI Studio", "Antigravity", "Gemini CLI"]
    roles = ["測點子 · 寫 prompt", "做出來 · Vibe Coding", "自動化 · 下指令"]

    for i in range(3):
        left = Inches(0.8 + i * 4.1)
        top = Inches(3.1)
        _card(s, left, top, Inches(3.7), Inches(3.0))
        _pill(s, titles[i], left + Inches(0.3), top + Inches(0.35), width=Inches(1.4), font_size=10)
        _add_text(s, names[i], left + Inches(0.3), top + Inches(1.05), Inches(3.2), Inches(0.7),
                  size=26, bold=True, letter_spacing=-0.8)
        _add_text(s, roles[i], left + Inches(0.3), top + Inches(1.85), Inches(3.2), Inches(0.6),
                  size=14, color=MUTED, letter_spacing=-0.1)
        _add_text(s, f"0{i + 1}", left + Inches(0.3), top + Inches(2.3), Inches(3.2), Inches(0.5),
                  size=14, color=MUTED, font=FONT_MONO, letter_spacing=3.0)

        if i < 2:
            arrow_left = left + Inches(3.85)
            _add_text(s, "→", arrow_left, top + Inches(1.3), Inches(0.25), Inches(0.6),
                      size=22, color=INK, align=PP_ALIGN.CENTER)

    _add_footer(s, "M0 · Opening", no)
    return s


def s_schedule_rules(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M0 · Schedule & rules", Inches(0.8), Inches(0.7))
    _add_text(s, "今日時程  ·  三個規則", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)

    schedule = [
        ("09:00", "M0", "開場"),
        ("09:30", "M1", "找痛點（你今天要做什麼）"),
        ("11:00", "M2", "STRIKE 戰法 × AI Studio"),
        ("13:30", "M3", "Antigravity Vibe Coding ⭐"),
        ("15:40", "M4", "Gemini CLI"),
        ("16:40", "M5", "部署 + 成果展示"),
    ]
    for i, (time, mod, topic) in enumerate(schedule):
        top = Inches(2.5 + i * 0.6)
        _hairline(s, Inches(0.8), top, Inches(6.8))
        _add_text(s, time, Inches(0.8), top + Inches(0.15), Inches(1.0), Inches(0.4),
                  size=14, color=INK, font=FONT_MONO, letter_spacing=1.0)
        _add_text(s, mod, Inches(1.85), top + Inches(0.15), Inches(0.8), Inches(0.4),
                  size=14, bold=True, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
        _add_text(s, topic, Inches(2.7), top + Inches(0.1), Inches(4.5), Inches(0.5), size=16, letter_spacing=-0.2)

    # 規則
    _hairline(s, Inches(8.2), Inches(2.5), Inches(4.3))
    _mono_label(s, "3 Rules", Inches(8.2), Inches(2.65))
    rules = [
        ("01", "不會的直接舉手"),
        ("02", "做錯沒關係，AI 陪你改"),
        ("03", "下課前能對朋友秀一個東西"),
    ]
    for i, (num, rule) in enumerate(rules):
        top = Inches(3.2 + i * 0.95)
        _add_text(s, num, Inches(8.2), top, Inches(0.8), Inches(0.5),
                  size=14, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
        _add_text(s, rule, Inches(8.2), top + Inches(0.35), Inches(4.3), Inches(0.6),
                  size=18, bold=True, letter_spacing=-0.2)

    _add_footer(s, "M0 · Opening", no)
    return s


# --- Shared templates ------------------------------------------------------

def s_module_divider(prs, no, module_code, module_name, subtitle, footer_text):
    """模組封面."""
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(3.6), height=Inches(0.15))

    _mono_label(s, f"Module {module_code}", Inches(0.8), Inches(1.2), size=12)
    _add_text(s, module_code, Inches(0.8), Inches(1.6), Inches(12), Inches(1.0),
              size=18, color=MUTED, font=FONT_MONO, letter_spacing=3.0)
    _add_text(s, module_name, Inches(0.8), Inches(2.1), Inches(12), Inches(1.4),
              size=88, bold=True, letter_spacing=-3.5)
    _add_text(s, subtitle, Inches(0.8), Inches(4.2), Inches(12), Inches(1.6),
              size=26, color=INK_SOFT, line_spacing=1.4, letter_spacing=-0.4)

    _add_footer(s, footer_text, no)
    return s


def s_practice_block(prs, no, module, title, steps, duration):
    """實作頁樣板（淡暖底 + 鮮明標示）."""
    s = _new_slide(prs, BG_SOFT)
    _pill(s, "Hands-on", Inches(0.8), Inches(0.7), width=Inches(1.3), fill=INK, text_color=BG, font_size=10)
    _add_text(s, title, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)
    _hairline(s, Inches(0.8), Inches(2.4), Inches(11.7), color=INK, weight=1.0)

    for i, step in enumerate(steps):
        _add_text(s, step, Inches(1.0), Inches(2.75 + i * 0.5), Inches(11.3), Inches(0.5),
                  size=18, color=INK_SOFT, letter_spacing=-0.2)

    _hairline(s, Inches(0.8), Inches(6.3), Inches(11.7))
    _add_text(s, "Duration", Inches(0.8), Inches(6.45), Inches(3), Inches(0.3),
              size=10, color=MUTED, font=FONT_MONO, letter_spacing=2.5)
    _add_text(s, duration, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
              size=22, bold=True, letter_spacing=-0.3)

    _add_footer(s, module, no)
    return s


def s_checkpoint(prs, no, cp_num, title, question, module):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.25))

    _mono_label(s, f"Checkpoint {cp_num:02d}", Inches(0.8), Inches(1.5))
    _add_text(s, f"CP{cp_num}", Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.5),
              size=72, bold=True, letter_spacing=-2.5)
    _add_text(s, title, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.8),
              size=24, color=INK_SOFT, letter_spacing=-0.3)
    _hairline(s, Inches(0.8), Inches(4.3), Inches(11.7), color=INK, weight=1.0)
    _add_text(s, question, Inches(0.8), Inches(4.55), Inches(11.7), Inches(1.5),
              size=26, bold=True, line_spacing=1.5, letter_spacing=-0.4)

    _add_text(s, "寫得出來 → 過關    ·    寫不出來 → 舉手，我們一起再拆一次",
              Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.0)
    _add_footer(s, module, no)
    return s


# --- M1 -------------------------------------------------------------------

def s_m1_cover(prs, no):
    return s_module_divider(
        prs, no, "M1", "找痛點",
        "做得出來，不等於做得對。\n選對題目，是今天最重要的決定。",
        "M1 · Pain Discovery",
    )


def s_m1_mistake(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M1 · The most common mistake", Inches(0.8), Inches(0.7))
    _add_text(s, "新手最容易犯的錯", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)

    blocks = [
        (Inches(0.8), "錯", "產品思維", "先想「我要做什麼產品」", "空想，容易越想越大\n做不完", ACCENT_WARN),
        (Inches(6.85), "對", "痛點思維", "先找「有什麼事讓我煩」", "具體、有人受苦\n做完就有用", ACCENT_GOOD),
    ]
    for left, mark, title, sub, desc, color in blocks:
        _card(s, left, Inches(2.5), Inches(5.7), Inches(3.9))
        _add_text(s, mark, left + Inches(0.35), Inches(2.7), Inches(0.8), Inches(0.8),
                  size=28, bold=True, color=color, letter_spacing=-0.5)
        _add_text(s, title, left + Inches(0.35), Inches(3.55), Inches(5), Inches(0.7),
                  size=28, bold=True, letter_spacing=-0.8)
        _hairline(s, left + Inches(0.35), Inches(4.3), Inches(5))
        _add_text(s, sub, left + Inches(0.35), Inches(4.45), Inches(5), Inches(0.6),
                  size=16, color=INK_SOFT, letter_spacing=-0.2)
        _add_text(s, desc, left + Inches(0.35), Inches(5.15), Inches(5), Inches(1.1),
                  size=15, color=MUTED, line_spacing=1.5, letter_spacing=-0.1)

    _add_footer(s, "M1 · Pain Discovery", no)
    return s


def s_m1_physics(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M1 · Four physics checks", Inches(0.8), Inches(0.7))
    _add_text(s, "一個好痛點的四個物理量", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)
    _add_text(s, "初學者第一個專案：① + ③ 過關就行。", Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=16, color=MUTED, letter_spacing=-0.2)

    items = [
        ("01", "具體的人", "能說出一個真人的名字（自己算）"),
        ("02", "正在付錢的痛", "他現在花錢、花時間、花精力在解決"),
        ("03", "手工可交付", "不寫程式，你今晚能幫他做一次"),
        ("04", "收款管道", "他願意付錢你收得到（之後再說）"),
    ]
    for i, (num, title, desc) in enumerate(items):
        row, col = divmod(i, 2)
        left = Inches(0.8 + col * 6.05)
        top = Inches(2.85 + row * 1.95)
        _card(s, left, top, Inches(5.7), Inches(1.7))
        _add_text(s, num, left + Inches(0.35), top + Inches(0.25), Inches(1.0), Inches(0.5),
                  size=13, color=MUTED, font=FONT_MONO, letter_spacing=2.5)
        _add_text(s, title, left + Inches(0.35), top + Inches(0.65), Inches(5), Inches(0.6),
                  size=22, bold=True, letter_spacing=-0.5)
        _add_text(s, desc, left + Inches(0.35), top + Inches(1.2), Inches(5), Inches(0.5),
                  size=14, color=MUTED, letter_spacing=-0.1)

    _add_footer(s, "M1 · Pain Discovery", no)
    return s


def s_m1_concrete(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M1 · What does concrete mean", Inches(0.8), Inches(0.7))
    _add_text(s, "範例：什麼叫「具體」", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)

    _mono_label(s, "× Vague", Inches(0.8), Inches(2.5), color=ACCENT_WARN)
    _mono_label(s, "✓ Concrete", Inches(6.85), Inches(2.5), color=ACCENT_GOOD)

    pairs = [
        ("年輕人", "我朋友阿明，28 歲，每週三晚上要做財報"),
        ("小企業老闆", "我認識的王太太，開早餐店，每天晚上手動記帳"),
        ("想學英文的人", "我表弟，下個月要考多益 750 分"),
    ]
    for i, (vague, concrete) in enumerate(pairs):
        top = Inches(3.1 + i * 1.0)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, vague, Inches(0.8), top + Inches(0.2), Inches(5.7), Inches(0.7),
                  size=18, color=MUTED, letter_spacing=-0.2)
        _add_text(s, concrete, Inches(6.85), top + Inches(0.2), Inches(5.7), Inches(0.7),
                  size=18, letter_spacing=-0.2)
    _hairline(s, Inches(0.8), Inches(6.1), Inches(11.7))
    _add_text(s, "痛點不抽象，不然沒辦法驗證。",
              Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
              size=20, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.3)

    _add_footer(s, "M1 · Pain Discovery", no)
    return s


def s_m1_prompt(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M1 · The magic prompt", Inches(0.8), Inches(0.7))
    _add_text(s, "找痛點的魔法 Prompt", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=36, bold=True, letter_spacing=-1.3)
    _add_text(s, "複製貼上到 ChatGPT / Claude", Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.4),
              size=14, color=MUTED, font=FONT_MONO, letter_spacing=1.5)

    code = """你是一位痛點結構化教練。我腦中有一些模糊的工作 / 生活痛點。

請用蘇格拉底式提問（一次一題）引導我整理出 3 個具體痛點，
每個痛點包含：
  - 誰遇到（具體的人）
  - 情境（什麼時候）
  - 現行解法（現在怎麼處理）
  - 不滿（哪裡煩）

規則：一次一題，我回答後再問下一題；
      我回答太抽象時請直接指出；
      不要評分或建議產品。

開始問我第一題吧。"""
    _add_code_block(s, code, Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.3), size=13)

    _add_footer(s, "M1 · Pain Discovery", no)
    return s


def s_m1_practice1(prs, no):
    return s_practice_block(
        prs, no, "M1 · Pain Discovery",
        "實作 1：找出你的 3 個痛點",
        [
            "01   打開 ChatGPT / Claude",
            "02   複製上一張的 Prompt",
            "03   跟它對話 15-20 分鐘",
            "04   產出 3 個具體痛點",
            "",
            "助教會巡場，卡住就舉手。",
        ],
        "25 分鐘",
    )


def s_m1_practice2(prs, no):
    s = _new_slide(prs, BG_SOFT)
    _pill(s, "Hands-on", Inches(0.8), Inches(0.7), width=Inches(1.3), fill=INK, text_color=BG, font_size=10)
    _add_text(s, "實作 2：物理量檢驗", Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)
    _add_text(s, "挑出你要做的那一個", Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2)

    code = """痛點 A   □ 具體的人    □ 付錢的痛    □ 手工可交付    □ 收款
痛點 B   □ 具體的人    □ 付錢的痛    □ 手工可交付    □ 收款
痛點 C   □ 具體的人    □ 付錢的痛    □ 手工可交付    □ 收款

選 ① + ③ 最紮實的那個  =  今天的題目"""
    _add_code_block(s, code, Inches(0.8), Inches(2.9), Inches(11.7), Inches(3.3), size=16)

    _hairline(s, Inches(0.8), Inches(6.3), Inches(11.7))
    _add_text(s, "Duration", Inches(0.8), Inches(6.45), Inches(3), Inches(0.3),
              size=10, color=MUTED, font=FONT_MONO, letter_spacing=2.5)
    _add_text(s, "15 分鐘", Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
              size=22, bold=True, letter_spacing=-0.3)

    _add_footer(s, "M1 · Pain Discovery", no)
    return s


def s_m1_one_sentence(prs, no):
    s = _new_slide(prs, BG_SOFT)
    _pill(s, "Hands-on", Inches(0.8), Inches(0.7), width=Inches(1.3), fill=INK, text_color=BG, font_size=10)
    _add_text(s, "把題目變成一句話", Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)
    _add_text(s, "用這個句型：", Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.4),
              size=14, color=MUTED, font=FONT_MONO, letter_spacing=1.5)

    code = """我要做 _______________，
給 _______________ 用，
解決 _______________ 問題，
核心功能是 _______________。

範例：
  我要做一個新聞摘要工具，
  給自己用，
  解決每天花 30 分鐘讀新聞只記得 10% 的問題，
  核心功能是貼新聞就出 3 句重點。"""
    _add_code_block(s, code, Inches(0.8), Inches(2.7), Inches(11.7), Inches(3.6), size=14)

    _hairline(s, Inches(0.8), Inches(6.4), Inches(11.7))
    _add_text(s, "15 分鐘  ·  下午 Antigravity 直接用這句話",
              Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.4),
              size=12, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)

    _add_footer(s, "M1 · Pain Discovery", no)
    return s


# --- M2 -------------------------------------------------------------------

def s_m2_cover(prs, no):
    return s_module_divider(
        prs, no, "M2", "STRIKE 提示詞戰法",
        "不是 AI 變聰明了。\n是你把需求講清楚了。",
        "M2 · STRIKE × AI Studio",
    )


def s_m2_compare(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M2 · Live comparison", Inches(0.8), Inches(0.7))
    _add_text(s, "現場對比：一句話  vs  完整 STRIKE",
              Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=34, bold=True, letter_spacing=-1.2)
    _add_text(s, "題目：會議記錄整理", Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.4),
              size=14, color=MUTED, font=FONT_MONO, letter_spacing=1.5)

    # A
    _mono_label(s, "A · One-liner", Inches(0.8), Inches(2.7), color=ACCENT_WARN)
    _add_code_block(s, '"幫我整理會議記錄"', Inches(0.8), Inches(3.1), Inches(5.85), Inches(0.7), size=13)

    # B
    _mono_label(s, "B · Full STRIKE", Inches(6.85), Inches(2.7), color=ACCENT_GOOD)
    strike = """【S】我是 PM，剛結束 Q2 跨部門規劃會議
【T】將逐字稿整理成結構化會議紀錄
【R】你是資深 PM，擅長提取關鍵決策
【I】[貼上 45 分鐘逐字稿]
【K】四區塊：摘要/決議/待辦(含負責人)/未解議題
       300-500 字
【E】## Q2 規劃會議  日期:..."""
    _add_code_block(s, strike, Inches(6.85), Inches(3.1), Inches(5.85), Inches(3.4), size=11)

    _add_text(s, "→ 現場開 AI Studio 跑，3 分鐘後比產出。這是今天第一個 AHA moment。",
              Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.3),
              size=12, color=MUTED, font=FONT_MONO, letter_spacing=1.0, align=PP_ALIGN.CENTER)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_ten_traps(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M2 · 10 common traps", Inches(0.8), Inches(0.7))
    _add_text(s, "你不是不會用 AI，", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=34, bold=True, letter_spacing=-1.2)
    _add_text(s, "是沒想清楚要什麼。", Inches(0.8), Inches(1.95), Inches(11.7), Inches(1.0),
              size=34, bold=True, letter_spacing=-1.2)

    traps = [
        "01  提示不夠精準", "06  缺少具體範例",
        "02  缺乏條理", "07  忽略限制條件",
        "03  缺乏背景資訊", "08  過於冗長或簡略",
        "04  目標不夠明確", "09  語言不一致",
        "05  提示過於模糊", "10  缺乏驗證與反饋",
    ]
    for i, t in enumerate(traps):
        row, col = divmod(i, 2)
        left = Inches(0.8 + col * 6.05)
        top = Inches(3.35 + row * 0.55)
        _add_text(s, t, left, top, Inches(5.7), Inches(0.45),
                  size=16, color=INK_SOFT, font=FONT_MONO, letter_spacing=0.5)

    _hairline(s, Inches(0.8), Inches(6.2), Inches(11.7), color=INK, weight=1.0)
    _add_text(s, "AI 產出品質  =  你定義問題的能力",
              Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.6),
              size=22, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.3)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_strike_hero(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M2 · The core framework", Inches(0.8), Inches(0.7))
    _add_text(s, "STRIKE 六字訣", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=44, bold=True, letter_spacing=-1.5)
    _add_text(s, "Prompt 工程的核心框架",
              Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.5),
              size=16, color=MUTED, letter_spacing=-0.2)

    items = [
        ("S", "Situation", "情境", "交代背景，讓 AI 不用猜"),
        ("T", "Task", "任務", "明確說做什麼、產出什麼"),
        ("R", "Role", "角色", "指定 AI 扮演的專業身分"),
        ("I", "Input", "輸入", "餵資料，不讓 AI 瞎猜"),
        ("K", "KPI", "驗收", "定義什麼叫「夠好」"),
        ("E", "Example", "範例", "給樣本，鎖定輸出方向"),
    ]
    for i, (letter, eng, cn, desc) in enumerate(items):
        row, col = divmod(i, 3)
        left = Inches(0.8 + col * 4.05)
        top = Inches(2.8 + row * 1.8)
        _card(s, left, top, Inches(3.85), Inches(1.55))
        _add_text(s, letter, left + Inches(0.25), top + Inches(0.05), Inches(1.0), Inches(1.2),
                  size=52, bold=True, letter_spacing=-2.5)
        _add_text(s, eng, left + Inches(1.4), top + Inches(0.25), Inches(2.4), Inches(0.35),
                  size=11, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
        _add_text(s, cn, left + Inches(1.4), top + Inches(0.55), Inches(2.4), Inches(0.45),
                  size=18, bold=True, letter_spacing=-0.3)
        _add_text(s, desc, left + Inches(1.4), top + Inches(1.0), Inches(2.4), Inches(0.5),
                  size=12, color=MUTED, letter_spacing=-0.1)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_strike_str(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M2 · Write well (1/2) — S / T / R", Inches(0.8), Inches(0.7))
    _add_text(s, "每格怎麼寫好（一）", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)

    blocks = [
        ("S  情境定場",
         '❌  空泛：「我想要...」',
         '✅  具體：「我是消費品行銷經理，準備 Q2 檢討會議」'),
        ("T  任務定向（動詞 + 產出物，一個 prompt 一件事）",
         '❌  「整理一下」',
         '✅  「整理成三段摘要，含待辦清單和負責人」'),
        ("R  角色定級（決定 AI 的專業天花板）",
         '❌  「你是助理」',
         '✅  「你是資深社群媒體分析師」'),
    ]
    for i, (title, bad, good) in enumerate(blocks):
        top = Inches(2.3 + i * 1.55)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, title, Inches(0.8), top + Inches(0.15), Inches(11.7), Inches(0.5),
                  size=18, bold=True, letter_spacing=-0.3)
        _add_text(s, bad, Inches(0.8), top + Inches(0.7), Inches(11.7), Inches(0.45),
                  size=13, color=ACCENT_WARN, font=FONT_MONO, letter_spacing=0.3)
        _add_text(s, good, Inches(0.8), top + Inches(1.1), Inches(11.7), Inches(0.45),
                  size=13, color=ACCENT_GOOD, font=FONT_MONO, letter_spacing=0.3)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_strike_ike(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M2 · Write well (2/2) — I / K / E", Inches(0.8), Inches(0.7))
    _add_text(s, "每格怎麼寫好（二）", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)

    blocks = [
        ("I  輸入定料（AI 沒資料就瞎猜）",
         '✅  「以下是我方 Q2 貼文數據：[貼上數據]」'),
        ("K  標準定規（定義什麼叫「夠好」）",
         "✅  格式：表格  |  字數：800 內  |  語氣：中立  |  禁止：推測"),
        ("E  範例定調（一個範例省三輪來回）",
         '✅  「參考這份過去的成功範本：[貼上]」'),
    ]
    for i, (title, good) in enumerate(blocks):
        top = Inches(2.3 + i * 1.2)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, title, Inches(0.8), top + Inches(0.15), Inches(11.7), Inches(0.5),
                  size=18, bold=True, letter_spacing=-0.3)
        _add_text(s, good, Inches(0.8), top + Inches(0.7), Inches(11.7), Inches(0.45),
                  size=13, color=ACCENT_GOOD, font=FONT_MONO, letter_spacing=0.3)

    _hairline(s, Inches(0.8), Inches(6.0), Inches(11.7), color=INK, weight=1.0)
    _add_text(s, "沒有 E 沒關係。沒有 T + K + I？那 AI 一定瞎掰。",
              Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8),
              size=22, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.3)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_mvp(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.25))
    _mono_label(s, "M2 · MVP formula", Inches(0.8), Inches(1.3))
    _add_text(s, "趕時間？救命組合。", Inches(0.8), Inches(1.75), Inches(11.7), Inches(0.8),
              size=26, color=MUTED, letter_spacing=-0.3)

    _add_text(s, "T  +  K  +  I", Inches(0.8), Inches(2.55), Inches(11.7), Inches(1.6),
              size=120, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-5.0)
    _add_text(s, "任務  ·  驗收  ·  輸入",
              Inches(0.8), Inches(4.25), Inches(11.7), Inches(0.6),
              size=22, color=INK_SOFT, align=PP_ALIGN.CENTER, letter_spacing=-0.3)

    _hairline(s, Inches(0.8), Inches(5.3), Inches(11.7))
    ratios = [("STRIKE 全寫", "100%", "滿分"), ("T + K + I", "80%", "可用"), ("都沒", "10%", "運氣")]
    for i, (name, pct, note) in enumerate(ratios):
        left = Inches(0.8 + i * 4.05)
        _add_text(s, name, left, Inches(5.5), Inches(3.85), Inches(0.4),
                  size=12, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)
        _add_text(s, pct, left, Inches(5.85), Inches(3.85), Inches(0.7),
                  size=36, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-1.2)
        _add_text(s, note, left, Inches(6.5), Inches(3.85), Inches(0.4),
                  size=13, color=MUTED, align=PP_ALIGN.CENTER, letter_spacing=-0.1)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_l123(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M2 · Three attack levels", Inches(0.8), Inches(0.7))
    _add_text(s, "三層進攻策略", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)
    _add_text(s, "配對今天三個工具。", Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=16, color=MUTED, letter_spacing=-0.2)

    levels = [
        ("L1", "試探射擊", "Recon Fire", "一句話丟出去，看 AI 能做什麼", "AI Studio 快速驗證"),
        ("L2", "精準打擊", "Precision Strike", "完整 STRIKE 六要素，一次到位", "PRD / Antigravity"),
        ("L3", "鏈式攻擊", "Chain Attack", "複雜任務拆 3-5 步，前一步 → 下一步", "Gemini CLI 管道 & 迭代"),
    ]
    for i, (lv, cn, en, desc, usage) in enumerate(levels):
        top = Inches(2.8 + i * 1.15)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, lv, Inches(0.8), top + Inches(0.15), Inches(1.2), Inches(0.6),
                  size=34, bold=True, letter_spacing=-1.0)
        _add_text(s, cn, Inches(2.2), top + Inches(0.25), Inches(3), Inches(0.5),
                  size=22, bold=True, letter_spacing=-0.4)
        _add_text(s, en, Inches(2.2), top + Inches(0.75), Inches(3), Inches(0.35),
                  size=11, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
        _add_text(s, desc, Inches(5.4), top + Inches(0.3), Inches(4.2), Inches(0.5),
                  size=14, color=INK_SOFT, letter_spacing=-0.1)
        _add_text(s, "→ " + usage, Inches(9.7), top + Inches(0.3), Inches(2.9), Inches(0.5),
                  size=12, color=MUTED, font=FONT_MONO, letter_spacing=1.0, align=PP_ALIGN.RIGHT)

    _hairline(s, Inches(0.8), Inches(6.3), Inches(11.7))
    _add_text(s, "三變數快篩：歧義度  ·  風險度  ·  自動化度  →  10 秒決定打法。",
              Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.5),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.0, align=PP_ALIGN.CENTER)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_ai_studio_intro(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M2 · The tool", Inches(0.8), Inches(0.7))
    _add_text(s, "Google AI Studio", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=44, bold=True, letter_spacing=-1.5)
    _add_text(s, "瀏覽器裡的 AI 實驗場。", Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2)

    features = [
        ("01", "零安裝", "免費，直接用 Gemini 最新模型"),
        ("02", "System Instruction", "放 S + R + K（設一次用多次）"),
        ("03", "Prompt 區", "每次放 T + I + E"),
        ("04", "Compare", "同一 prompt 比兩個模型 / 兩種寫法"),
    ]
    for i, (num, title, desc) in enumerate(features):
        top = Inches(2.9 + i * 0.75)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, num, Inches(0.8), top + Inches(0.2), Inches(0.8), Inches(0.5),
                  size=12, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
        _add_text(s, title, Inches(1.8), top + Inches(0.15), Inches(4), Inches(0.5),
                  size=18, bold=True, letter_spacing=-0.3)
        _add_text(s, desc, Inches(6.2), top + Inches(0.2), Inches(6.3), Inches(0.5),
                  size=15, color=INK_SOFT, letter_spacing=-0.2)

    _hairline(s, Inches(0.8), Inches(6.25), Inches(11.7))
    _add_text(s, "aistudio.google.com",
              Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
              size=18, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER, letter_spacing=1.0)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_practice_strike(prs, no):
    s = _new_slide(prs, BG_SOFT)
    _pill(s, "Hands-on", Inches(0.8), Inches(0.7), width=Inches(1.3), fill=INK, text_color=BG, font_size=10)
    _add_text(s, "實作 1：STRIKE 填空套版", Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)

    code = """【S】我是 _____________ ，正在 _____________ 。
【T】請 [動詞] 一份 [產出物]
【R】你是 _____________（專業身分）
【I】[貼上你的資料 / 痛點描述]
【K】格式：___ ；字數：___ ；禁止：___
【E】參考格式：___

在 AI Studio 測 3 次，每次只改一格，看產出變化。"""
    _add_code_block(s, code, Inches(0.8), Inches(2.4), Inches(11.7), Inches(3.7), size=15)

    _hairline(s, Inches(0.8), Inches(6.25), Inches(11.7))
    _add_text(s, "25 分鐘  ·  關鍵：一次只改一格",
              Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_practice_prd(prs, no):
    s = _new_slide(prs, BG_SOFT)
    _pill(s, "Hands-on", Inches(0.8), Inches(0.7), width=Inches(1.3), fill=INK, text_color=BG, font_size=10)
    _add_text(s, "實作 2：用 STRIKE 寫出你的 PRD", Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0),
              size=34, bold=True, letter_spacing=-1.2)
    _add_text(s, "PRD 本身就是一份 L2 精準打擊。",
              Inches(0.8), Inches(2.25), Inches(11.7), Inches(0.5),
              size=16, color=MUTED, letter_spacing=-0.2)

    code = """跟 AI 說：
"基於上面 STRIKE，請幫我寫一份產品規格文件（PRD）：
  ① 產品一句話（S + T）
  ② 目標使用者（S）
  ③ 功能清單（T 拆成 5 項以內）
  ④ 輸入 / 輸出範例（I + E）
  ⑤ 驗收標準（K：介面、技術限制、不做的事）"

把產出複製到共用文件，下午丟進 Antigravity。"""
    _add_code_block(s, code, Inches(0.8), Inches(2.85), Inches(11.7), Inches(3.4), size=13)

    _hairline(s, Inches(0.8), Inches(6.4), Inches(11.7))
    _add_text(s, "20 分鐘  ·  砍得狠，下午才做得成",
              Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.5),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


def s_m2_defense(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M2 · Defense & recap", Inches(0.8), Inches(0.7))
    _add_text(s, "AI 會一本正經地胡說八道。", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=30, bold=True, letter_spacing=-1.0)
    _add_text(s, "三句話救你。", Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.8),
              size=30, bold=True, letter_spacing=-1.0)

    defenses = [
        ('「為什麼？」', "逼它交代推理"),
        ('「還有其他方案嗎？」', "打破第一直覺"),
        ('「幫我找出這段的漏洞」', "讓 AI 自攻自己"),
    ]
    for i, (q, purpose) in enumerate(defenses):
        left = Inches(0.8 + i * 4.05)
        top = Inches(3.1)
        _card(s, left, top, Inches(3.85), Inches(1.5))
        _add_text(s, q, left + Inches(0.25), top + Inches(0.35), Inches(3.5), Inches(0.6),
                  size=17, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.3)
        _add_text(s, purpose, left + Inches(0.25), top + Inches(0.95), Inches(3.5), Inches(0.4),
                  size=12, color=MUTED, align=PP_ALIGN.CENTER, letter_spacing=-0.1)

    _hairline(s, Inches(0.8), Inches(5.0), Inches(11.7), color=INK)
    _mono_label(s, "M2 Recap", Inches(0.8), Inches(5.15))
    recap = [
        "STRIKE 六字訣：S / T / R / I / K / E",
        "救命組合：T + K + I（80% 可用）",
        "打法選擇：L1 試探 · L2 精準 · L3 鏈式",
        "你手上有一份 PRD，下午變成可運行的東西。",
    ]
    for i, line in enumerate(recap):
        _add_text(s, "·  " + line, Inches(0.8), Inches(5.55 + i * 0.35), Inches(11.7), Inches(0.35),
                  size=13, color=INK_SOFT, letter_spacing=-0.1)

    _add_footer(s, "M2 · STRIKE × AI Studio", no)
    return s


# --- M3 -------------------------------------------------------------------

def s_m3_cover(prs, no):
    return s_module_divider(
        prs, no, "M3", "Antigravity · Vibe Coding",
        "你描述需求  ·  AI 負責寫程式碼。\n這是今天的主場。",
        "M3 · Antigravity",
    )


def s_m3_what_is(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M3 · Definition", Inches(0.8), Inches(0.7))
    _add_text(s, "什麼是 Vibe Coding？", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)
    _add_text(s, "不寫程式碼，用自然語言描述你要什麼。AI 把你的想法變成可運行的產品。",
              Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.8),
              size=18, color=MUTED, letter_spacing=-0.2)

    _hairline(s, Inches(0.8), Inches(3.2), Inches(11.7), color=INK)

    _mono_label(s, "✓ Your job", Inches(0.8), Inches(3.35), color=ACCENT_GOOD)
    _mono_label(s, "× Not your job", Inches(6.85), Inches(3.35), color=ACCENT_WARN)

    yours = ["定義問題", "描述需求", "驗證結果", "修正方向"]
    theirs = ["寫程式碼", "記語法", "除 bug", "看 stack trace"]
    for i, (a, b) in enumerate(zip(yours, theirs)):
        top = Inches(3.9 + i * 0.55)
        _add_text(s, a, Inches(0.8), top, Inches(5.7), Inches(0.5), size=20, bold=True, letter_spacing=-0.3)
        _add_text(s, b, Inches(6.85), top, Inches(5.7), Inches(0.5), size=20, color=MUTED, letter_spacing=-0.3)

    _add_footer(s, "M3 · Antigravity", no)
    return s


def s_m3_antigravity_intro(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M3 · The tool", Inches(0.8), Inches(0.7))
    _add_text(s, "Antigravity", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=44, bold=True, letter_spacing=-1.5)
    _add_text(s, "Agent-first IDE。不是自動補全，是整個專案都幫你做。",
              Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2)

    pairs = [
        ("你是任務管理者", "AI 是開發者"),
        ("你描述需求", "AI 規劃 + 寫程式碼"),
        ("你看到每行改動", "AI 解釋為什麼這樣寫"),
    ]
    for i, (left, right) in enumerate(pairs):
        top = Inches(3.1 + i * 0.85)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, left, Inches(0.8), top + Inches(0.2), Inches(5.8), Inches(0.5),
                  size=20, bold=True, letter_spacing=-0.3)
        _add_text(s, right, Inches(6.85), top + Inches(0.2), Inches(5.8), Inches(0.5),
                  size=20, color=INK_SOFT, letter_spacing=-0.3)

    _hairline(s, Inches(0.8), Inches(6.0), Inches(11.7))
    _add_text(s, "Free · Public Preview  ·  antigravity.google",
              Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
              size=15, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)

    _add_footer(s, "M3 · Antigravity", no)
    return s


def s_m3_five_steps(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M3 · The loop", Inches(0.8), Inches(0.7))
    _add_text(s, "Vibe Coding 五步流程", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)

    steps = [
        ("01", "描述需求", "用 M2 的 STRIKE（PRD = L2 精準打擊）"),
        ("02", "看 AI 規劃", "AI 列出：做什麼檔案、用什麼技術"),
        ("03", "看程式碼", "不用懂，但要看。不懂就問 AI"),
        ("04", "測試", "跑起來、點一點、貼資料試"),
        ("05", "自然語言修正", "「太擠」「加匯出」「顏色改深藍」"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        top = Inches(2.4 + i * 0.75)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, num, Inches(0.8), top + Inches(0.2), Inches(1.0), Inches(0.5),
                  size=16, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
        _add_text(s, title, Inches(2.0), top + Inches(0.15), Inches(3.5), Inches(0.5),
                  size=22, bold=True, letter_spacing=-0.4)
        _add_text(s, desc, Inches(5.8), top + Inches(0.2), Inches(6.8), Inches(0.5),
                  size=14, color=MUTED, letter_spacing=-0.2)

    _add_text(s, "↻  回到 03、04 直到滿意。",
              Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)

    _add_footer(s, "M3 · Antigravity", no)
    return s


def s_m3_golden_quote(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.25))
    _mono_label(s, "The most important mindset", Inches(0.8), Inches(1.3))

    _add_text(s, "AI 做錯了，", Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.0),
              size=50, color=MUTED, letter_spacing=-1.5, align=PP_ALIGN.CENTER)
    _add_text(s, "不是去改 code，", Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.3),
              size=80, bold=True, letter_spacing=-3.0, align=PP_ALIGN.CENTER)
    _add_text(s, "是去改描述。", Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.3),
              size=80, bold=True, letter_spacing=-3.0, align=PP_ALIGN.CENTER)

    _add_footer(s, "M3 · Antigravity", no)
    return s


def s_m3_natural_lang_fix(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M3 · How to fix", Inches(0.8), Inches(0.7))
    _add_text(s, "自然語言修正範例", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)

    _mono_label(s, "× Don't", Inches(0.8), Inches(2.5), color=ACCENT_WARN)
    _mono_label(s, "✓ Do this", Inches(6.85), Inches(2.5), color=ACCENT_GOOD)

    pairs = [
        ("打開 .css 改色", '「首頁主色改成深藍 #1E40AF」'),
        ("去修 JS 邏輯", '「搜尋功能沒動作，請檢查」'),
        ("自己調排版", '「卡片間距太擠，間隔加大 30%」'),
        ("複製 error 問它", '「跑起來白畫面，請幫我找問題」'),
    ]
    for i, (bad, good) in enumerate(pairs):
        top = Inches(3.1 + i * 0.7)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, bad, Inches(0.8), top + Inches(0.18), Inches(5.7), Inches(0.5),
                  size=16, color=MUTED, letter_spacing=-0.2)
        _add_text(s, good, Inches(6.85), top + Inches(0.18), Inches(5.7), Inches(0.5),
                  size=16, font=FONT_MONO, letter_spacing=0.3)
    _hairline(s, Inches(0.8), Inches(5.9), Inches(11.7), color=INK)
    _add_text(s, "修正方向  ≠  修改程式碼",
              Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8),
              size=24, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.4)

    _add_footer(s, "M3 · Antigravity", no)
    return s


def s_m3_demo_time(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.25))
    _mono_label(s, "Live demo — 15 min", Inches(0.8), Inches(1.3))

    _add_text(s, "Demo Time", Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5),
              size=96, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-4.0)
    _add_text(s, "講師把一份 PRD 丟進 Antigravity，全程走完五步流程，現場做出一個小工具。",
              Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.8),
              size=18, color=INK_SOFT, align=PP_ALIGN.CENTER, letter_spacing=-0.2)

    _hairline(s, Inches(2.5), Inches(5.3), Inches(8.3))
    _add_text(s, "學員記下：", Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
              size=14, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)
    notes = ["AI 怎麼規劃", "AI 的 code 大概長什麼樣", "講師怎麼「說話」修正"]
    for i, n in enumerate(notes):
        _add_text(s, "·  " + n, Inches(0.8), Inches(5.95 + i * 0.35), Inches(11.7), Inches(0.35),
                  size=14, color=INK_SOFT, align=PP_ALIGN.CENTER, letter_spacing=-0.1)

    _add_footer(s, "M3 · Antigravity", no)
    return s


def s_m3_practice1(prs, no):
    return s_practice_block(
        prs, no, "M3 · Antigravity",
        "實作：你的專案，你來做（第一輪）",
        [
            "01   打開 Antigravity",
            "02   建新專案",
            "03   把你的 PRD 貼到第一個 prompt",
            "04   讓 AI 規劃 → 看 code → 測試 → 修正",
            "05   跑到一個「可以看的版本」",
            "",
            "卡關 2 分鐘以上：舉手。",
        ],
        "60 分鐘",
    )


def s_m3_common_issues(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M3 · Common issues", Inches(0.8), Inches(0.7))
    _add_text(s, "常見問題  ×  處理", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)

    _mono_label(s, "Situation", Inches(0.8), Inches(2.5))
    _mono_label(s, "Response", Inches(6.85), Inches(2.5))

    issues = [
        ("AI 產出的 UI 超醜", '「介面改得更簡潔現代，用 Tailwind」'),
        ("功能少了一塊", '「請加上 XX 功能，規格同 PRD 第 N 項」'),
        ("跑起來錯誤訊息", "把訊息貼給 AI，請它修"),
        ("AI 一直繞圈圈", "重開新 chat，重新描述"),
        ("AI 做太多、加了沒要的", '「請砍掉 XX 功能，我只要 YY」'),
    ]
    for i, (a, b) in enumerate(issues):
        top = Inches(3.1 + i * 0.65)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, a, Inches(0.8), top + Inches(0.15), Inches(5.7), Inches(0.5),
                  size=16, color=INK_SOFT, letter_spacing=-0.2)
        _add_text(s, b, Inches(6.85), top + Inches(0.15), Inches(5.7), Inches(0.5),
                  size=14, font=FONT_MONO, letter_spacing=0.3)

    _add_footer(s, "M3 · Antigravity", no)
    return s


def s_m3_practice2(prs, no):
    return s_practice_block(
        prs, no, "M3 · Antigravity",
        "實作：加一個新功能（第二輪）",
        [
            "選一個（或你自己想的）：",
            "□   加匯出 CSV / JSON",
            "□   加深色主題切換",
            "□   加歷史紀錄",
            "□   加搜尋 / 篩選",
            "",
            "完全用自然語言跟 Antigravity 說。迭代才是 Vibe Coding 的精髓。",
        ],
        "20 分鐘",
    )


def s_m3_cp2(prs, no):
    return s_checkpoint(
        prs, no, 2,
        "打開你的電腦，你能看到……",
        "□  你的專案跑起來了\n□  可以貼資料 / 點按鈕試功能\n□  經過至少一次自然語言修正",
        "M3 · Antigravity",
    )


def s_m3_summary(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M3 · Recap", Inches(0.8), Inches(0.7))
    _add_text(s, "M3 小結", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)

    items = [
        "Vibe Coding = 描述需求比寫 code 重要",
        "AI 做錯 → 改描述，不是改 code",
        "迭代才是精髓，一次成功只是運氣",
        "你已經有一個可運行的小專案了",
    ]
    for i, line in enumerate(items):
        top = Inches(2.5 + i * 0.7)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, "·  " + line, Inches(0.8), top + Inches(0.15), Inches(11.7), Inches(0.55),
                  size=22, letter_spacing=-0.3)
    _hairline(s, Inches(0.8), Inches(5.4), Inches(11.7))

    _add_text(s, "下一站：我們讓它變得更強 →",
              Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
              size=16, color=INK_SOFT, align=PP_ALIGN.CENTER, letter_spacing=-0.2)
    _add_text(s, "☕  休息 10 分鐘",
              Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
              size=12, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)

    _add_footer(s, "M3 · Antigravity", no)
    return s


# --- M4 -------------------------------------------------------------------

def s_m4_cover(prs, no):
    return s_module_divider(
        prs, no, "M4", "Gemini CLI",
        "從「介面操作」到「打字就走」。\n順便學 2026 通用 coding agent 技能。",
        "M4 · Gemini CLI",
    )


def s_m4_terminal(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · The terminal is friendly", Inches(0.8), Inches(0.7))
    _add_text(s, "終端機沒那麼可怕", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)
    _add_text(s, "終端機  =  用打字代替點滑鼠。",
              Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2)

    _add_text(s, "你已經會：", Inches(0.8), Inches(2.9), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5)
    _add_code_block(s, "打開網頁 → 輸入網址 → Enter", Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.6), size=15)

    _add_text(s, "跟這個一樣，只是：", Inches(0.8), Inches(4.15), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5)
    _add_code_block(s, "打開終端 → 輸入指令 → Enter", Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.6), size=15)

    _add_text(s, "三個指令你馬上就會：", Inches(0.8), Inches(5.35), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5)
    code = 'gemini "hello"\ngemini "把這句翻成英文：今天天氣真好"\ncat news.txt | gemini "用三句話摘要"'
    _add_code_block(s, code, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.1), size=13)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_continuity(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · Continuity across tools", Inches(0.8), Inches(0.7))
    _add_text(s, "你在 M2 學的 STRIKE，", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=36, bold=True, letter_spacing=-1.2)
    _add_text(s, "可以直接當 CLI 指令用。", Inches(0.8), Inches(1.95), Inches(11.7), Inches(1.0),
              size=36, bold=True, letter_spacing=-1.2)

    mappings = [
        ("貼到 AI Studio", "瀏覽器看結果"),
        ("丟進 Antigravity", "IDE 裡做專案"),
        ("打進 Gemini CLI", "終端機立刻跑"),
    ]
    for i, (a, b) in enumerate(mappings):
        top = Inches(3.3 + i * 0.85)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, a, Inches(0.8), top + Inches(0.2), Inches(5.7), Inches(0.5),
                  size=22, bold=True, align=PP_ALIGN.RIGHT, letter_spacing=-0.4)
        _add_text(s, "→", Inches(6.5), top + Inches(0.2), Inches(0.5), Inches(0.5),
                  size=22, color=MUTED, align=PP_ALIGN.CENTER)
        _add_text(s, b, Inches(7.1), top + Inches(0.2), Inches(5.5), Inches(0.5),
                  size=22, color=INK_SOFT, letter_spacing=-0.4)
    _hairline(s, Inches(0.8), Inches(6.0), Inches(11.7), color=INK)
    _add_text(s, "差別只在介面，不在 prompt。",
              Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6),
              size=22, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.4)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_commands(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · Five commands to use today", Inches(0.8), Inches(0.7))
    _add_text(s, "五個你馬上能用的指令",
              Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=36, bold=True, letter_spacing=-1.2)

    code = """# 1. 直接問
gemini "幫我把這句中文翻成英文：你好嗎"

# 2. 管道：把檔案丟給 AI
cat news.txt | gemini "用三句話摘要"

# 3. 寫入輸出
gemini "這段 CSV 轉成 JSON" < data.csv > data.json

# 4. 套 prompt 模板
gemini "你是編輯(R)，用五要格式(K)摘要：$(cat input.txt)"

# 5. L3 鏈式（前一步 | 後一步）
cat news.txt | gemini "摘要" | gemini "翻成英文"

# 進階：互動模式
gemini"""
    _add_code_block(s, code, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.6), size=13)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_three_abilities(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · The shape of a coding agent", Inches(0.8), Inches(0.7))
    _add_text(s, "認識你的 coding agent", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)
    _add_text(s, "你今天碰到的三個工具，其實都是同一種東西。",
              Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2)

    abilities = [
        ("🗣  聊天",    "跟你對話、寫程式碼"),
        ("📂  看檔",    "讀你電腦上的檔案"),
        ("🔧  用工具",  "搜尋、執行指令、呼叫服務"),
    ]
    for i, (a, b) in enumerate(abilities):
        top = Inches(3.1 + i * 0.85)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, a, Inches(0.8), top + Inches(0.2), Inches(4.5), Inches(0.5),
                  size=24, bold=True, letter_spacing=-0.4)
        _add_text(s, b, Inches(5.5), top + Inches(0.2), Inches(7.0), Inches(0.5),
                  size=20, color=INK_SOFT, letter_spacing=-0.3)
    _hairline(s, Inches(0.8), Inches(5.85), Inches(11.7), color=INK)

    _add_text(s, "差別只在「介面」與「能用哪些工具」。",
              Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.6),
              size=22, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.4)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_gemini_md(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · GEMINI.md — project memory", Inches(0.8), Inches(0.7))
    _add_text(s, "GEMINI.md", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, font=FONT_MONO, letter_spacing=-1.0)
    _add_text(s, "給 agent 一份「公司守則」。每次開工自動讀。",
              Inches(0.8), Inches(2.05), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2)

    code = """# 我的新聞摘要工具

- 語言：Python
- 永遠用繁體中文回覆
- 摘要請限制在 100 字內
- 程式碼請加註解"""
    _add_code_block(s, code, Inches(0.8), Inches(2.85), Inches(11.7), Inches(2.2), size=14)

    _add_text(s, "兩個層級（Gemini CLI 都吃）：",
              Inches(0.8), Inches(5.25), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5)
    levels = """~/.gemini/GEMINI.md   = 全域（所有專案都套用）
./GEMINI.md           = 這個專案專用"""
    _add_code_block(s, levels, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.1), size=14)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_at_ref(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · @ reference — feed files into chat", Inches(0.8), Inches(0.7))
    _add_text(s, "@ 引用：把檔案丟進對話",
              Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=36, bold=True, letter_spacing=-1.2)

    _add_text(s, "舊招（M4 前半教過）：",
              Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5)
    _add_code_block(s, 'cat news.txt | gemini "摘要"',
                    Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.6), size=15)

    _add_text(s, "新招（更簡單）：",
              Inches(0.8), Inches(3.55), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5)
    _add_code_block(s, 'gemini "@news.txt 幫我摘要"',
                    Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.6), size=15)

    _hairline(s, Inches(0.8), Inches(4.85), Inches(11.7))
    _add_text(s, "cat |  =  把檔案內容「灌」進去（單檔、純文字）",
              Inches(0.8), Inches(5.05), Inches(11.7), Inches(0.5),
              size=16, color=INK_SOFT, letter_spacing=-0.2)
    _add_text(s, "@file  =  告訴 agent「自己去讀」（多檔、大檔、含程式碼）",
              Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.5),
              size=16, color=INK_SOFT, letter_spacing=-0.2)

    _add_code_block(s, 'gemini "@README.md @src/main.py 解釋這個程式怎麼跑"',
                    Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6), size=14)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_tools_perm(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · Built-in tools & safety", Inches(0.8), Inches(0.7))
    _add_text(s, "內建工具 + 安全模式",
              Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=36, bold=True, letter_spacing=-1.2)

    _add_text(s, "Gemini CLI 不只會聊天，預設帶這些工具：",
              Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=16, color=MUTED, letter_spacing=-0.2)

    tools = [
        ("📂  ReadFile / WriteFile", "讀寫檔案"),
        ("🐚  Shell",                "跑指令"),
        ("🔍  GoogleSearch",         "即時上網查"),
        ("🌐  WebFetch",             "讀網頁"),
    ]
    for i, (a, b) in enumerate(tools):
        top = Inches(2.9 + i * 0.5)
        _add_text(s, a, Inches(1.0), top, Inches(5.5), Inches(0.4),
                  size=18, bold=True, font=FONT_MONO, letter_spacing=-0.2)
        _add_text(s, b, Inches(6.8), top, Inches(5.7), Inches(0.4),
                  size=18, color=INK_SOFT, letter_spacing=-0.2)

    _hairline(s, Inches(0.8), Inches(5.05), Inches(11.7), color=INK)
    _add_text(s, "每次它要動手，會先問你「可以嗎？」",
              Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.5),
              size=18, bold=True, letter_spacing=-0.3)
    _add_text(s, "預設模式：每步都確認  ← 初學者請用這個",
              Inches(1.0), Inches(5.85), Inches(11.7), Inches(0.4),
              size=15, color=ACCENT_GOOD, letter_spacing=-0.2)
    _add_text(s, "--yolo 模式：不問直接做  ← demo 用，別在重要專案開",
              Inches(1.0), Inches(6.3), Inches(11.7), Inches(0.4),
              size=15, color=ACCENT_WARN, font=FONT_MONO, letter_spacing=-0.2)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_mcp(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · MCP — the USB-C for AI", Inches(0.8), Inches(0.7))
    _add_text(s, "MCP：AI 的 USB-C",
              Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)

    _add_text(s, "問題：每個工具都要重新接一次 Slack、GitHub、資料庫，超累。",
              Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.5),
              size=16, color=MUTED, letter_spacing=-0.2)
    _add_text(s, "解法：Model Context Protocol —— 業界共通的「插槽標準」。",
              Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.5),
              size=18, bold=True, letter_spacing=-0.3)

    _add_text(s, "常見 MCP server（會了就有得用）：",
              Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5)
    mcps = [
        ("🔧  filesystem",  "讓 agent 管理你的檔案"),
        ("🌐  fetch",       "讓 agent 抓網頁"),
        ("📅  google-cal",  "讓 agent 看你的行事曆"),
        ("💬  slack",       "讓 agent 收發訊息"),
    ]
    for i, (a, b) in enumerate(mcps):
        top = Inches(4.0 + i * 0.5)
        _add_text(s, a, Inches(1.0), top, Inches(5.0), Inches(0.4),
                  size=17, bold=True, font=FONT_MONO, letter_spacing=-0.2)
        _add_text(s, b, Inches(6.3), top, Inches(6.2), Inches(0.4),
                  size=17, color=INK_SOFT, letter_spacing=-0.2)

    _hairline(s, Inches(0.8), Inches(6.2), Inches(11.7), color=INK)
    _add_text(s, "今天不裝。先知道有這東西，回家想玩再裝。",
              Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
              size=18, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.3)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_cross_tools(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · Same skill, different names", Inches(0.8), Inches(0.7))
    _add_text(s, "你學的不是 Gemini CLI，是通用技能",
              Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0),
              size=32, bold=True, letter_spacing=-1.2)

    header = "                AI Studio   Antigravity      Gemini CLI    Claude Code   Cursor"
    rows = """專案記憶        ❌          AGENTS.md *      GEMINI.md     CLAUDE.md     .cursor/rules/
@ 檔案引用      ❌          ✅              ✅            ✅            ✅
STRIKE          ✅          ✅              ✅            ✅            ✅
MCP             ❌          ✅              ✅            ✅            ✅"""
    _add_code_block(s, header + "\n" + rows,
                    Inches(0.4), Inches(2.3), Inches(12.5), Inches(2.4), size=12)

    _add_text(s, "*  Antigravity 還會讀 ~/.gemini/GEMINI.md 做全域設定（跟 Gemini CLI 共用——同一家人）",
              Inches(0.8), Inches(4.95), Inches(11.7), Inches(0.4),
              size=13, color=MUTED, letter_spacing=-0.1)

    _hairline(s, Inches(0.8), Inches(5.55), Inches(11.7), color=INK)
    _add_text(s, "🟡  AGENTS.md 是 2026 新興的「跨工具通用」格式",
              Inches(0.8), Inches(5.75), Inches(11.7), Inches(0.5),
              size=20, bold=True, letter_spacing=-0.3)
    _add_text(s, "    Cursor 與 Antigravity 都認得它，最值得學",
              Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.5),
              size=18, color=INK_SOFT, letter_spacing=-0.2)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_practice(prs, no):
    return s_practice_block(
        prs, no, "M4 · Gemini CLI",
        "實作：為你的專案加一點自動化",
        [
            "必做（5 min）：在你的專案根目錄建一份 GEMINI.md（5 行就好）",
            "",
            "三選一（20 min）：",
            "A.   把你專案的核心功能寫成一行指令",
            "     例：summarize.sh news.txt",
            "",
            "B.   寫一個自動產測試資料的腳本",
            "     例：gen_data.sh 10  → 產 10 筆假資料",
            "",
            "C.   寫一個「每天跑一次」的小工具（不實際排程）",
        ],
        "25 分鐘",
    )


def s_m4_decision_tree(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · Tool decision tree", Inches(0.8), Inches(0.7))
    _add_text(s, "三工具決策樹", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)
    _add_text(s, "我想做一件事，選什麼工具？",
              Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=16, color=MUTED, letter_spacing=-0.2)

    tree = """├─  在瀏覽器、想測想法               →   AI Studio
├─  做有介面的東西                   →   Antigravity
├─  把重複的事自動化                 →   Gemini CLI
├─  超大專案 / 多人協作              →   Antigravity（多 agent）
└─  馬上用、不想安裝                 →   AI Studio"""
    _add_code_block(s, tree, Inches(1.2), Inches(2.9), Inches(10.9), Inches(2.6), size=15)

    _hairline(s, Inches(0.8), Inches(5.9), Inches(11.7))
    _add_text(s, '請 2-3 位學員分享：「我今天三個工具用下來，最驚喜的是 ___」',
              Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8),
              size=18, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.3)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


def s_m4_summary(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M4 · Recap + CP3", Inches(0.8), Inches(0.7))
    _add_text(s, "M4 小結", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)

    items = [
        "終端機 = 打字的瀏覽器",
        "Gemini CLI = AI Studio 的無介面版",
        "同一套 STRIKE，三工具通用；CLI 最適合 L3 鏈式",
        "CLI 是你之後做「每天自動跑」的基礎",
    ]
    for i, line in enumerate(items):
        top = Inches(2.4 + i * 0.55)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, "·  " + line, Inches(0.8), top + Inches(0.12), Inches(11.7), Inches(0.45),
                  size=18, letter_spacing=-0.2)
    _hairline(s, Inches(0.8), Inches(4.6), Inches(11.7), color=INK)

    _mono_label(s, "Checkpoint 03", Inches(0.8), Inches(4.8))
    _add_text(s, "能對旁邊的人口述：這三個工具我各自會用來做什麼。",
              Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.0),
              size=24, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.4)

    _add_footer(s, "M4 · Gemini CLI", no)
    return s


# --- M5 -------------------------------------------------------------------

def s_m5_cover(prs, no):
    return s_module_divider(
        prs, no, "M5", "最後一哩路",
        "部署  ·  成果展示  ·  行動計畫。\n做出來不算數，能打開網址才算。",
        "M5 · Deploy & Wrap-up",
    )


def s_m5_deploy(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M5 · Three-step deploy", Inches(0.8), Inches(0.7))
    _add_text(s, "快速部署三步驟", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=40, bold=True, letter_spacing=-1.3)
    _add_text(s, "把你的 Antigravity 專案丟上網。",
              Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=16, color=MUTED, letter_spacing=-0.2)

    steps = [
        ("01", "程式碼 push 到 GitHub", "Antigravity 有內建 Git 按鈕"),
        ("02", "去 Cloudflare Pages / Vercel", "登入 → Import Git Repo → 選你的 repo"),
        ("03", "自動部署 → 取得公開網址", "三步驟搞定"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        top = Inches(2.9 + i * 1.1)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, num, Inches(0.8), top + Inches(0.25), Inches(1.0), Inches(0.5),
                  size=28, bold=True, font=FONT_MONO, letter_spacing=-0.5)
        _add_text(s, title, Inches(2.3), top + Inches(0.2), Inches(10), Inches(0.5),
                  size=22, bold=True, letter_spacing=-0.4)
        _add_text(s, desc, Inches(2.3), top + Inches(0.75), Inches(10), Inches(0.45),
                  size=14, color=MUTED, letter_spacing=-0.2)

    _add_text(s, "講師現場示範一次（15 min）",
              Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.4),
              size=12, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)

    _add_footer(s, "M5 · Deploy & Wrap-up", no)
    return s


def s_m5_practice(prs, no):
    return s_practice_block(
        prs, no, "M5 · Deploy & Wrap-up",
        "自由時間：選一個",
        [
            "□   把專案部署上線（跟著講師做）",
            "□   打磨你的專案（加最後一個小功能）",
            "□   寫你的 30 天行動卡",
            "□   跟旁邊的人 demo 彼此的專案",
            "",
            "助教隨時支援。",
        ],
        "15 分鐘",
    )


def s_m5_show(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.25))
    _mono_label(s, "Showcase · 15 min", Inches(0.8), Inches(1.3))
    _add_text(s, "成果展示", Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.2),
              size=72, bold=True, letter_spacing=-2.5)
    _add_text(s, "自願者上台 3 分鐘。",
              Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2)

    items = [
        ("01", "你做了什麼？", "1 min"),
        ("02", "最大的 AHA moment 是哪個？", "1 min"),
        ("03", "下一步你想做什麼？", "1 min"),
    ]
    for i, (num, q, t) in enumerate(items):
        top = Inches(4.0 + i * 0.7)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, num, Inches(0.8), top + Inches(0.2), Inches(0.9), Inches(0.5),
                  size=14, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
        _add_text(s, q, Inches(1.95), top + Inches(0.15), Inches(8.5), Inches(0.5),
                  size=20, letter_spacing=-0.3)
        _add_text(s, t, Inches(10.5), top + Inches(0.2), Inches(2), Inches(0.4),
                  size=12, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.RIGHT)

    _add_footer(s, "M5 · Deploy & Wrap-up", no)
    return s


def s_m5_action_card(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M5 · 30-day action card", Inches(0.8), Inches(0.7))
    _add_text(s, "30 天行動卡", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)
    _add_text(s, "回去之後，每週做一件事。",
              Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.5),
              size=14, color=MUTED, letter_spacing=-0.2)

    weeks = [
        ("Week 01", "工具內化",
         ["今天的專案加一個新功能", "用 STRIKE 六字訣跑一次全新任務", "用 Gemini CLI 做一件以前手動的事"]),
        ("Week 02", "新專案啟動",
         ["找一個新痛點（跑物理量四問）", "用 AI Studio 寫新專案 PRD", "在 Antigravity 做出原型"]),
        ("Week 03", "分享 & 回饋",
         ["把一個專案上線", "找 1 位朋友用、收 3 點回饋", "依回饋改一個功能"]),
        ("Week 04", "沉澱",
         ["寫筆記：三工具我各自什麼時候用", "寫一個「每天自動跑」CLI 腳本", "想：下個月我想做什麼？"]),
    ]
    for i, (wk, title, items) in enumerate(weeks):
        row, col = divmod(i, 2)
        left = Inches(0.8 + col * 6.05)
        top = Inches(2.85 + row * 2.0)
        _card(s, left, top, Inches(5.7), Inches(1.85))
        _mono_label(s, wk, left + Inches(0.3), top + Inches(0.2), color=MUTED)
        _add_text(s, title, left + Inches(0.3), top + Inches(0.5), Inches(5.0), Inches(0.45),
                  size=17, bold=True, letter_spacing=-0.3)
        for j, item in enumerate(items):
            _add_text(s, "□  " + item, left + Inches(0.3), top + Inches(0.95 + j * 0.3), Inches(5.2), Inches(0.3),
                      size=11, color=INK_SOFT, letter_spacing=-0.1)

    _add_footer(s, "M5 · Deploy & Wrap-up", no)
    return s


def s_m5_takeaway(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "M5 · What you're taking home", Inches(0.8), Inches(0.7))
    _add_text(s, "今天你帶走了什麼", Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=38, bold=True, letter_spacing=-1.3)

    items = [
        ("01", "一個可運行的小專案（選修已上線）"),
        ("02", "你的 PRD（存在 Notion / HackMD）"),
        ("03", "STRIKE 速查卡 + 追問三指令小卡"),
        ("04", "工具選擇決策樹"),
        ("05", "痛點物理量檢驗單"),
        ("06", "30 天行動卡"),
    ]
    for i, (num, desc) in enumerate(items):
        top = Inches(2.4 + i * 0.55)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _add_text(s, num, Inches(0.8), top + Inches(0.12), Inches(1), Inches(0.5),
                  size=13, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
        _add_text(s, desc, Inches(2.0), top + Inches(0.1), Inches(10.5), Inches(0.5),
                  size=18, letter_spacing=-0.3)

    _hairline(s, Inches(0.8), Inches(6.0), Inches(11.7), color=INK)
    _add_text(s, "+  最重要的：你知道自己做得到。",
              Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6),
              size=22, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.4)

    _add_footer(s, "M5 · Deploy & Wrap-up", no)
    return s


def s_m5_golden_closing(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.25))
    _mono_label(s, "Closing", Inches(0.8), Inches(1.2))

    _add_text(s, "AI 時代的贏家，", Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8),
              size=28, color=MUTED, letter_spacing=-0.5, align=PP_ALIGN.CENTER)
    _add_text(s, "不是最會寫程式的人，",
              Inches(0.8), Inches(2.55), Inches(11.7), Inches(0.8),
              size=34, letter_spacing=-1.0, align=PP_ALIGN.CENTER)
    _add_text(s, "不是最會用工具的人。",
              Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.8),
              size=34, letter_spacing=-1.0, align=PP_ALIGN.CENTER)

    _hairline(s, Inches(3.5), Inches(4.3), Inches(6.3), color=INK)
    _add_text(s, "而是", Inches(0.8), Inches(4.45), Inches(11.7), Inches(0.5),
              size=18, color=MUTED, letter_spacing=-0.2, align=PP_ALIGN.CENTER)

    _add_text(s, "最會定義問題的人。",
              Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.0),
              size=50, bold=True, letter_spacing=-1.8, align=PP_ALIGN.CENTER)
    _add_text(s, "能讓 AI 幫你解題的人。",
              Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0),
              size=50, bold=True, letter_spacing=-1.8, align=PP_ALIGN.CENTER)
    _add_text(s, "── 你今天已經跨出那一步。",
              Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
              size=12, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)

    _add_footer(s, "M5 · Deploy & Wrap-up", no)
    return s


def s_m5_thanks(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(0), height=Inches(0.25))
    _hero_gradient_strip(s, top=Inches(7.25), height=Inches(0.25))

    _mono_label(s, "Thank you", Inches(0.8), Inches(1.3))
    _add_text(s, "感謝你今天的一天。",
              Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.2),
              size=64, bold=True, letter_spacing=-2.5, align=PP_ALIGN.CENTER)

    _hairline(s, Inches(3), Inches(3.7), Inches(7.3))
    _mono_label(s, "Next steps", Inches(0.8), Inches(3.9), width=Inches(11.7), size=12)

    lines = [
        "[ 你的 email ]",
        "[ 學員社群 LINE / Discord ]",
        "[ 課程網頁 / 你的 Notion ]",
    ]
    for i, line in enumerate(lines):
        _add_text(s, line, Inches(0.8), Inches(4.4 + i * 0.45), Inches(11.7), Inches(0.4),
                  size=16, font=FONT_MONO, align=PP_ALIGN.CENTER, letter_spacing=1.0)

    _hairline(s, Inches(3), Inches(6.0), Inches(7.3))
    _add_text(s, "下一場工作坊：Vibe Coding 進階  ·  部署實戰  ·  商業化",
              Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5),
              size=13, color=MUTED, font=FONT_MONO, letter_spacing=1.5, align=PP_ALIGN.CENTER)
    _add_text(s, "合照時間  📸",
              Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.4),
              size=14, color=INK_SOFT, letter_spacing=-0.2, align=PP_ALIGN.CENTER)

    return s


# --- Appendix ---

def s_appendix_cover(prs, no):
    s = _new_slide(prs, BG)
    _hero_gradient_strip(s, top=Inches(3.6), height=Inches(0.15))
    _mono_label(s, "Appendix", Inches(0.8), Inches(1.2), size=12)
    _add_text(s, "附錄", Inches(0.8), Inches(1.6), Inches(12), Inches(1.0),
              size=18, color=MUTED, font=FONT_MONO, letter_spacing=3.0)
    _add_text(s, "三個備選 Demo 專案",
              Inches(0.8), Inches(2.1), Inches(12), Inches(1.4),
              size=72, bold=True, letter_spacing=-2.5)
    _add_text(s, "講師 M3 現場 Demo 三選一  ·  學員選題亦可參考。",
              Inches(0.8), Inches(4.2), Inches(12), Inches(0.6),
              size=20, color=INK_SOFT, letter_spacing=-0.2)
    _add_text(s, "共同特徵：免費額度充裕  ·  介面單頁  ·  2 小時內可完成原型。",
              Inches(0.8), Inches(4.9), Inches(12), Inches(0.6),
              size=16, color=MUTED, letter_spacing=-0.1)

    _add_footer(s, "Appendix · Demo candidates", no)
    return s


def _demo_slide(prs, no, tag, title, subtitle, ai_studio, antigravity, cli, why):
    s = _new_slide(prs, BG)
    _mono_label(s, tag, Inches(0.8), Inches(0.7))
    _add_text(s, title, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=34, bold=True, letter_spacing=-1.2)
    _add_text(s, subtitle, Inches(0.8), Inches(2.05), Inches(11.7), Inches(0.5),
              size=15, color=MUTED, letter_spacing=-0.2)

    stages = [
        ("AI Studio", ai_studio),
        ("Antigravity", antigravity),
        ("Gemini CLI", cli),
    ]
    for i, (name, desc) in enumerate(stages):
        top = Inches(2.75 + i * 0.85)
        _hairline(s, Inches(0.8), top, Inches(11.7))
        _pill(s, name, Inches(0.8), top + Inches(0.18), width=Inches(1.9), fill=INK, text_color=BG, font_size=10)
        _add_text(s, desc, Inches(3.0), top + Inches(0.18), Inches(9.5), Inches(0.6),
                  size=14, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE, letter_spacing=-0.1)

    _hairline(s, Inches(0.8), Inches(5.35), Inches(11.7), color=INK)
    _mono_label(s, "Why recommend", Inches(0.8), Inches(5.5))
    _add_text(s, why, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.1),
              size=14, color=INK_SOFT, line_spacing=1.5, letter_spacing=-0.1)

    _add_footer(s, "Appendix · Demo candidates", no)
    return s


def s_demo_news(prs, no):
    return _demo_slide(
        prs, no,
        "Demo A · News summarizer",
        "每日新聞 / 電子報 AI 摘要助手",
        "銜接性最佳的首選（教案原生示範）",
        "寫 STRIKE prompt：新聞全文 → 3 點重點 + 關注角度 + 偏見提醒",
        "單頁網頁：貼新聞 → 按鈕 → 摘要卡片（可匯出 Markdown）",
        'cat news.txt | gemini "摘要" | gemini "翻成英文"（L3 鏈式）',
        "教案 S16、S33 現成材料；Gemini 免費純文字額度綽綽有餘；L3 鏈式天然銜接 M4 CLI；30 人同時跑無壓力。",
    )


def s_demo_meeting(prs, no):
    return _demo_slide(
        prs, no,
        "Demo B · Meeting minutes",
        "會議逐字稿 → 結構化會議紀錄",
        "上班族 AHA 最強、專業感最高",
        "測 prompt：逐字稿 → 摘要 / 決議 / 待辦（含負責人）/ 未解議題",
        "網頁：貼逐字稿 → 產 Markdown 會議紀錄，一鍵複製到 Notion",
        "批次處理多場會議：for f in meetings/*.txt; do gemini < $f; done",
        "產出物可直接在公司用；對應 M2 [S16] 的對比範例；AHA 情緒曲線最強；只需準備 1-2 份逐字稿素材。",
    )


def s_demo_pantry(prs, no):
    return _demo_slide(
        prs, no,
        "Demo C · Pantry chef",
        "冰箱食材 → 今晚煮什麼",
        "最生活化、最輕鬆，適合下午氣氛放鬆時",
        "測 prompt：食材清單 + 口味偏好 → 3 道菜建議 + 步驟 + 缺什麼",
        "網頁：食材 checkbox + 時間限制 + 口味 → 食譜卡片",
        'gemini "今晚煮什麼" --pantry fridge.txt（可定時每天 17:00 跑）',
        "情境輕鬆，學員放鬆；現場可問「你冰箱有什麼？」互動感強；token 消耗極低；最貼近生活第 3 點需求。",
    )


def s_appendix_summary(prs, no):
    s = _new_slide(prs, BG)
    _mono_label(s, "Appendix · Comparison", Inches(0.8), Inches(0.7))
    _add_text(s, "備選 Demo 比較速查",
              Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.0),
              size=34, bold=True, letter_spacing=-1.2)

    headers = ["專案", "專業性", "生活化", "免費額度", "銜接性", "建議用途"]
    rows = [
        ["A  新聞摘要", "★★★", "★★★★", "★★★★★", "★★★★★", "講師主 Demo（推薦）"],
        ["B  會議紀錄", "★★★★★", "★★★", "★★★★★", "★★★★", "學員選題（上班族）"],
        ["C  冰箱食材", "★★", "★★★★★", "★★★★★", "★★★★", "下午氣氛放鬆備案"],
    ]
    col_widths = [2.4, 1.5, 1.5, 1.6, 1.5, 3.2]
    col_lefts = [0.8]
    for w in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + w)

    # 表頭
    _hairline(s, Inches(0.8), Inches(2.4), Inches(11.7), color=INK)
    for i, h in enumerate(headers):
        _add_text(s, h, Inches(col_lefts[i]), Inches(2.5), Inches(col_widths[i]), Inches(0.4),
                  size=11, color=MUTED, font=FONT_MONO, letter_spacing=2.0)
    _hairline(s, Inches(0.8), Inches(3.0), Inches(11.7))

    for ri, row in enumerate(rows):
        top = Inches(3.2 + ri * 1.0)
        for ci, cell in enumerate(row):
            size = 16 if ci == 0 else 15
            bold = ci == 0
            _add_text(s, cell, Inches(col_lefts[ci]), top + Inches(0.2), Inches(col_widths[ci]), Inches(0.5),
                      size=size, bold=bold, letter_spacing=-0.2)
        _hairline(s, Inches(0.8), top + Inches(0.85), Inches(11.7))

    _add_text(s, "我的建議：主 Demo 用 A 新聞摘要，備案 C 冰箱食材應付氣氛變化。",
              Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
              size=16, bold=True, align=PP_ALIGN.CENTER, letter_spacing=-0.3)

    _add_footer(s, "Appendix · Demo candidates", no)
    return s


# --- Builder ---------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        # M0 (5)
        s_cover, s_you_too, s_promise, s_three_tools, s_schedule_rules,
        # M1 (9)
        s_m1_cover, s_m1_mistake, s_m1_physics, s_m1_concrete, s_m1_prompt,
        s_m1_practice1, s_m1_practice2, s_m1_one_sentence,
        lambda p, n: s_checkpoint(p, n, 1, "你能寫出這句話嗎？",
                                  "「我今天要做 ___，給 ___ 用，\n解決 ___ 問題。」",
                                  "M1 · Pain Discovery"),
        # M2 (12)
        s_m2_cover, s_m2_compare, s_m2_ten_traps, s_m2_strike_hero,
        s_m2_strike_str, s_m2_strike_ike, s_m2_mvp, s_m2_l123,
        s_m2_ai_studio_intro, s_m2_practice_strike, s_m2_practice_prd, s_m2_defense,
        # M3 (12)
        s_m3_cover, s_m3_what_is, s_m3_antigravity_intro, s_m3_five_steps,
        s_m3_golden_quote, s_m3_natural_lang_fix, s_m3_demo_time,
        s_m3_practice1, s_m3_common_issues, s_m3_practice2,
        s_m3_cp2, s_m3_summary,
        # M4 (13)
        s_m4_cover, s_m4_terminal, s_m4_continuity, s_m4_commands,
        s_m4_three_abilities, s_m4_gemini_md, s_m4_at_ref,
        s_m4_tools_perm, s_m4_mcp, s_m4_cross_tools,
        s_m4_practice, s_m4_decision_tree, s_m4_summary,
        # M5 (8)
        s_m5_cover, s_m5_deploy, s_m5_practice, s_m5_show,
        s_m5_action_card, s_m5_takeaway, s_m5_golden_closing, s_m5_thanks,
        # Appendix (5)
        s_appendix_cover, s_demo_news, s_demo_meeting, s_demo_pantry, s_appendix_summary,
    ]

    for idx, b in enumerate(builders, start=1):
        b(prs, idx)

    prs.save(OUT_PATH)
    print(f"[OK] saved: {OUT_PATH}")
    print(f"[OK] total slides: {len(builders)}")


if __name__ == "__main__":
    build()
