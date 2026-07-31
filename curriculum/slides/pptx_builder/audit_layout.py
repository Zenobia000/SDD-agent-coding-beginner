"""Vibe Coding 投影片排版稽核器.

讀產出的 vibe_coding_workshop.pptx，找出**視覺上會有問題的 slide**。
跑：`uv run python audit_layout.py`

偵測重點（只報嚴重問題、忽略設計上的小重疊）：
- 文字「實際渲染」會超出 6.95"（撞 footer hairline）
- shape 右邊界超出投影片寬度 + 0.5"
- 內文字級 < 12pt 且不是 chrome（footer / mono label）
- shape 重疊面積 > 0.4 平方吋（設計上不會有的大塊重疊）

退出碼：嚴重違規 > 0 時為 1。
"""

from __future__ import annotations

import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.shapes.connector import Connector
from pptx.util import Emu, Inches, Pt

PPTX_PATH = Path(__file__).resolve().parent.parent / "vibe_coding_workshop.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SAFE_BOTTOM = Inches(6.95)
RIGHT_TOL = Inches(0.5)
MIN_BODY_FONT_PT = 12
MIN_CHROME_FONT_PT = 9
CHROME_TOP_THRESHOLD = Inches(6.2)
LINE_SPACING = 1.35
OVERLAP_AREA_TOL = Inches(0.4) * Inches(0.4)


def _in(emu: int | None) -> float:
    return (emu or 0) / 914400


def _fmt(emu: int | None) -> str:
    return f"{_in(emu):.2f}\""


def _shape_box(shape) -> tuple[int, int, int, int]:
    return (shape.left or 0, shape.top or 0,
            (shape.left or 0) + (shape.width or 0),
            (shape.top or 0) + (shape.height or 0))


def _overlap_area(a, b) -> int:
    ox = max(0, min(a[2], b[2]) - max(a[0], b[0]))
    oy = max(0, min(a[3], b[3]) - max(a[1], b[1]))
    return ox * oy


def _is_background_rect(shape) -> bool:
    if not shape.width or not shape.height:
        return False
    return shape.width >= SLIDE_W * 0.99 and shape.height >= SLIDE_H * 0.99


def _rendered_text_height(shape) -> int:
    """估算實際渲染高度（不是 textbox 框高度）."""
    if not shape.has_text_frame:
        return 0
    total_lines = 0
    max_pt = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text:
                if run.font.size is not None:
                    max_pt = max(max_pt, run.font.size.pt)
        total_lines += 1
    if max_pt == 0 or total_lines == 0:
        return 0
    pt_to_emu = 914400 / 72
    return int(max_pt * LINE_SPACING * total_lines * pt_to_emu)


def _is_chrome(shape) -> bool:
    """頁尾、頁碼、mono caps label——這些字級小是設計刻意的."""
    if not shape.has_text_frame:
        return False
    if (shape.top or 0) > CHROME_TOP_THRESHOLD:
        return True
    text = shape.text_frame.text
    return text.isupper() and len(text) < 60 and ("·" in text or text.strip().count(" ") <= 3)


def _is_pill_text(shape) -> bool:
    """pill 內部的文字（小、置中）不該被當成需要 ≥12pt 的內文."""
    if not shape.has_text_frame:
        return False
    if shape.height and shape.height < Inches(0.5) and shape.width and shape.width < Inches(2.0):
        return True
    return False


def _min_font_in_shape(shape):
    if not shape.has_text_frame:
        return None
    sizes = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text and run.font.size is not None:
                sizes.append(run.font.size.pt)
    return min(sizes) if sizes else None


def audit():
    if not PPTX_PATH.exists():
        print(f"[ERROR] 找不到 {PPTX_PATH}，請先跑 build_slides.py")
        sys.exit(2)

    prs = Presentation(PPTX_PATH)
    violations = defaultdict(list)
    total_slides = len(prs.slides)

    for idx, slide in enumerate(prs.slides, start=1):
        shapes_with_text = []

        for shape in slide.shapes:
            if _is_background_rect(shape) or isinstance(shape, Connector):
                continue

            left, top = shape.left or 0, shape.top or 0
            width, height = shape.width or 0, shape.height or 0
            right = left + width

            # 1. 渲染文字撞 footer 或落在 footer 下方（非 chrome 字級）
            if shape.has_text_frame and shape.text_frame.text.strip():
                rendered_h = _rendered_text_height(shape)
                rendered_bottom = top + rendered_h
                min_size = _min_font_in_shape(shape)
                is_chrome_size = min_size is not None and min_size <= MIN_CHROME_FONT_PT + 1
                if rendered_bottom > SAFE_BOTTOM and not is_chrome_size:
                    if top >= SAFE_BOTTOM:
                        violations[idx].append(
                            f"[HIGH] 文字落在 footer 之下：top={_fmt(top)}（>{_fmt(SAFE_BOTTOM)}），"
                            f"字級 {min_size}pt 應為內文 文字「{shape.text_frame.text[:25]}…」"
                        )
                    else:
                        violations[idx].append(
                            f"[HIGH] 文字渲染撞 footer：top={_fmt(top)} + 渲染高={_fmt(rendered_h)} "
                            f"→ 底部={_fmt(rendered_bottom)}（安全 ≤{_fmt(SAFE_BOTTOM)}）"
                            f" 文字「{shape.text_frame.text[:25]}…」"
                        )

            # 2. 右邊界超出
            if right > SLIDE_W + RIGHT_TOL:
                snippet = shape.text_frame.text[:25] if shape.has_text_frame else "(非文字)"
                violations[idx].append(
                    f"[HIGH] 右邊界超出：right={_fmt(right)}（投影片寬 {_fmt(SLIDE_W)}）"
                    f" shape「{snippet}…」"
                )

            # 3. 字級過小（區分 chrome 與內文）
            min_size = _min_font_in_shape(shape)
            if min_size is not None:
                threshold = MIN_CHROME_FONT_PT if (_is_chrome(shape) or _is_pill_text(shape)) else MIN_BODY_FONT_PT
                if min_size < threshold:
                    snippet = shape.text_frame.text[:25]
                    violations[idx].append(
                        f"[MED] 字級過小：{min_size:.0f}pt（要求 ≥{threshold}pt）"
                        f" 文字「{snippet}…」"
                    )

            # 收集 box 給後續重疊檢測
            if shape.has_text_frame and shape.text_frame.text.strip():
                shapes_with_text.append((shape, _shape_box(shape)))

        # 4. 大塊 shape 重疊
        for i in range(len(shapes_with_text)):
            for j in range(i + 1, len(shapes_with_text)):
                area = _overlap_area(shapes_with_text[i][1], shapes_with_text[j][1])
                if area > OVERLAP_AREA_TOL:
                    t1 = shapes_with_text[i][0].text_frame.text[:18]
                    t2 = shapes_with_text[j][0].text_frame.text[:18]
                    overlap_in = (area ** 0.5) / 914400
                    violations[idx].append(
                        f"[MED] 大塊重疊：「{t1}」╳「{t2}」 約 {overlap_in:.2f}\"² "
                    )

    print(f"\n{'=' * 64}")
    print(f"  Vibe Coding 投影片排版稽核  ({total_slides} 張)")
    print(f"{'=' * 64}\n")

    if not violations:
        print("✅ 全部通過。\n")
        return 0

    high_count = sum(
        1 for msgs in violations.values() if any("[HIGH]" in m for m in msgs)
    )
    print(f"✅ 通過：{total_slides - len(violations)} / {total_slides}")
    print(f"❌ 違規：{len(violations)} 張（其中 {high_count} 張含 HIGH 風險）\n")

    for slide_no in sorted(violations.keys()):
        print(f"  [S{slide_no:02d}]")
        for msg in violations[slide_no]:
            print(f"    · {msg}")
        print()

    return 1 if high_count > 0 else 0


if __name__ == "__main__":
    sys.exit(audit())
